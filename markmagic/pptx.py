"""Convert pptx to markdown."""

import re
from typing import IO

import pptx
from pptx.shapes.autoshape import Shape
from pptx.shapes.graphfrm import GraphicFrame
from tabulate import tabulate

from markmagic.settings import Settings


def convert_pptx(file: str | IO[bytes], settings: Settings) -> str:
    """Convert pptx to markdown."""
    markdown = ""
    presentation = pptx.Presentation(file)
    slide_num = 0
    for slide in presentation.slides:
        slide_num += 1
        markdown += f"\n\n<!-- Slide number: {slide_num} -->\n\n"
        title_shape = slide.shapes.title
        for shape in slide.shapes:
            if shape.has_table:
                if isinstance(shape, GraphicFrame):
                    markdown += _convert_table(shape, settings)
            elif shape.has_text_frame:
                if isinstance(shape, Shape):
                    if shape == title_shape:
                        markdown += "# " + shape.text.lstrip() + "\n"
                    else:
                        markdown += shape.text + "\n"
            # TODO: Consider something for pictures
        if slide.has_notes_slide:
            markdown += "\n### Notes:\n"
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame is not None:
                markdown += notes_frame.text
        markdown = markdown.strip()
    return markdown


def _convert_table(graphfrm: GraphicFrame, settings: Settings) -> str:
    tabular_data = [
        [re.sub(r"\s+", " ", cell.text).strip() for cell in row.cells] for row in graphfrm.table.rows
    ]
    markdown = (
        tabulate(
            tabular_data,
            tablefmt=settings.tablefmt,
            showindex=settings.showindex,
            headers="firstrow",
        )
        + "\n\n"
    )
    return markdown
